# Файл: src/converters/pptx_engine.py
"""
Движок конвертации презентаций PowerPoint (.pptx) в структурированный Markdown.
"""

from pathlib import Path
from typing import Dict, Any, List
from pptx import Presentation

from src.converters.base import BaseConverter, ConversionResult


class PptxConverter(BaseConverter):
    """
    Конвертер PowerPoint (.pptx) -> Markdown.
    """

    def can_handle(self, file_path: Path) -> bool:
        return file_path.suffix.lower() in [".pptx", ".pptm", ".potx"]

    def convert(self, file_path: Path, extract_profile: bool = False) -> ConversionResult:
        if not file_path.exists():
            raise FileNotFoundError(f"Файл презентации '{file_path}' не найден!")

        prs = Presentation(str(file_path))
        slides_count = len(prs.slides)
        md_sections: List[str] = [f"# Презентация: {file_path.name}\n"]
        total_tables = 0

        for idx, slide in enumerate(prs.slides, 1):
            title = f"Слайд {idx}"
            slide_text_blocks: List[str] = []
            notes_text = ""

            # Ищем заголовок слайда и текстовые блоки
            for shape in slide.shapes:
                if shape.has_text_frame:
                    tf = shape.text_frame
                    paragraphs = [p.text.strip() for p in tf.paragraphs if p.text.strip()]
                    if not paragraphs:
                        continue

                    # Если это заголовок слайда
                    if shape == slide.shapes.title or (idx == 1 and not slide_text_blocks):
                        title = f"Слайд {idx}: {paragraphs[0]}"
                        if len(paragraphs) > 1:
                            slide_text_blocks.extend([f"- {p}" for p in paragraphs[1:]])
                    else:
                        for p in tf.paragraphs:
                            p_text = p.text.strip()
                            if not p_text:
                                continue
                            indent_level = p.level or 0
                            prefix = "  " * indent_level + "- "
                            slide_text_blocks.append(f"{prefix}{p_text}")

                # Таблицы на слайдах
                elif shape.has_table:
                    total_tables += 1
                    t = shape.table
                    table_rows = []
                    for row in t.rows:
                        row_vals = [c.text.strip().replace("\n", " ") for c in row.cells]
                        table_rows.append(row_vals)
                    if table_rows:
                        cols = len(table_rows[0])
                        slide_text_blocks.append("\n| " + " | ".join(table_rows[0]) + " |")
                        slide_text_blocks.append("| " + " | ".join([":---"] * cols) + " |")
                        for r in table_rows[1:]:
                            padded = r + [""] * (cols - len(r))
                            slide_text_blocks.append("| " + " | ".join(padded[:cols]) + " |")
                        slide_text_blocks.append("\n")

            # Заметки спикера / докладчика
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    notes_text = f"\n> 💡 **Заметки докладчика:** {notes}\n"

            md_sections.append(f"## {title}\n")
            if slide_text_blocks:
                md_sections.append("\n".join(slide_text_blocks))
            if notes_text:
                md_sections.append(notes_text)
            md_sections.append("\n---\n")

        full_md = "\n".join(md_sections)

        layout_report = (
            f"📽️ ПРЕЗЕНТАЦИЯ POWERPOINT ({file_path.name}):\n"
            f"   • Количество слайдов: {slides_count}\n"
            f"   • Таблиц на слайдах: {total_tables}"
        )

        return ConversionResult(
            markdown=full_md,
            profile=None,
            metadata={
                "source_file": str(file_path),
                "slides_count": slides_count,
                "tables_count": total_tables,
                "format": "PPTX",
            },
            layout_report=layout_report,
            tables_count=total_tables,
            pages_count=slides_count,
        )
